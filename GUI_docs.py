import openpyxl as xl
import os, time
import keyboard as ky
import pyperclip
import tkinter as tk
from functools import partial
from pptx import Presentation



root=tk.Tk()
root.title("Create Files")
v = tk.IntVar()
languages = [
    ("Excel"),
    ("PPT"),
    ("Word"),
    ("New mail")
]

def create_xl(e1,e3):
    wb = xl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    file_n=e3.get()+".xlsx"
    b=e1.get()
    loc=b.replace("\\","\\\\")
    #global fn
    fn=loc+file_n
    wb.save(filename = fn)
    #print(wb.sheetnames)
    #file_path=os.path.abspath(fn)
    print(fn)
    ky.press_and_release('win+r')    
    time.sleep(1)
    pyperclip.copy(fn)
    ky.press_and_release('ctrl+v')
    ky.press_and_release('enter')

def create_ppt(e1,e3):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    #title = slide.shapes.title
    #subtitle = slide.placeholders[1]
    file_n=e3.get()+".pptx"
    b=e1.get()
    loc=b.replace("\\","\\\\")
    #global fn
    fn=loc+file_n
    #title.text = "Hello, World!"
    #subtitle.text = "python-pptx was here!"
    prs.save(fn)
    #file_path=os.path.abspath(fn)
    print(fn)
    ky.press_and_release('win+r')    
    time.sleep(1)
    pyperclip.copy(fn)
    ky.press_and_release('ctrl+v')
    ky.press_and_release('enter')


def XL():
    master = tk.Toplevel(root) 
    master.title("EXCEL")
    master.grab_set()
    tk.Label(master, text="Directory").grid(row=0)
    tk.Label(master, text="Name").grid(row=1)
    e1 = tk.Entry(master)
    e3 = tk.Entry(master)
    e1.grid(row=0, column=1)
    e3.grid(row=1, column=1)
    tk.Button(master, text='Close', command=master.destroy).grid(row=3, column=0,  pady=4)
    tk.Button(master, text='Create', command=partial(create_xl,e1,e3)).grid(row=3, column=1,  pady=4)
    #tk.Button(master, text='Open', command=open_xl).grid(row=3, column=2,  pady=4)    
    #master.mainloop( )

def PPT():
    master = tk.Toplevel(root) 
    master.title("PPT")
    master.grab_set()
    tk.Label(master, text="Directory").grid(row=0)
    tk.Label(master, text="Name").grid(row=1)
    e1 = tk.Entry(master)
    e3 = tk.Entry(master)
    e1.grid(row=0, column=1)
    e3.grid(row=1, column=1)
    tk.Button(master, text='Close', command=master.destroy).grid(row=3, column=0,  pady=4)
    tk.Button(master, text='Create', command=partial(create_ppt,e1,e3)).grid(row=3, column=1,  pady=4)
    #tk.Button(master, text='Open', command=open_xl).grid(row=3, column=2,  pady=4)    
    #master.mainloop( )
    
def ShowChoice():
    if (v.get()==0):
        XL()
    if(v.get()==1):
        PPT()
    if(v.get()==2):
        DOC()
    if(v.get()==3):
        MAIL()


tk.Label(root, 
         text="Choose the action to be performed",
         justify = tk.LEFT,
         padx = 20).pack()

for val, language in enumerate(languages):
    tk.Radiobutton(root,text=language,padx = 20,variable=v,command=ShowChoice,value=val).pack(anchor=tk.W)
tk.Button(root, text='Close', command=root.destroy).pack()
root.mainloop()
